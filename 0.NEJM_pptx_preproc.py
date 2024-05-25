import os
from pptx import Presentation
import pandas as pd

def extract_text_from_slide(slide):
    """
    Extract text from text boxes in a slide.
    Skip text boxes with content "Image Challenge" or "Q:".
    """
    text_boxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_content = text_frame.text.strip()
        if text_content in ["Image Challenge", "Q:"]:
            continue
        text_boxes.append(text_content)
    return text_boxes

def extract_images_from_pptx(pptx_path):
    """
    Extract images from a PPTX file and save them in a folder.
    Also extract text from text boxes and save both in an Excel file.
    """
    prs = Presentation(pptx_path)

    # Create 'pptimages' folder if it doesn't exist
    output_folder = 'pptimages'
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    data = []  # List to store data for Excel file

    for slide_number, slide in enumerate(prs.slides, start=1):
        # Extract text from slide
        slide_texts = extract_text_from_slide(slide)

        img_index = 0  # Initialize image index for the slide
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Check if shape is a picture
                image = shape.image
                image_width, image_height = image.size
                # Save image only if it's larger than 300x60
                if image_width > 300 and image_height > 60:
                    image_bytes = image.blob
                    image_filename = f'{output_folder}/img_page{slide_number}_{img_index}'
                    with open(f'{image_filename}.png', 'wb') as img_file:
                        img_file.write(image_bytes)
                    # Add image filename and associated text to data list
                    for text in slide_texts:
                        data.append([image_filename, text])
                    img_index += 1

    # Save data to Excel file
    save_text_to_excel(data, 'NEJMImageChallenge.xlsx')

    print(f'Images are extracted and saved in {output_folder}, skipping smaller images.')
    print('Texts are extracted and saved in NEJMImageChallenge.xlsx.')

def save_text_to_excel(data, excel_path):
    """
    Save extracted text and image paths to an Excel file.
    """
    df = pd.DataFrame(data, columns=["Image path", "NEJM-Image-Challenge"])
    df.to_excel(excel_path, index=False)

def main():
    # Specify the path to the "COMBINED" first page(s) download PPTX file
    pptx_path = 'NEJMImageChallenge_Combined.pptx'
    extract_images_from_pptx(pptx_path)

if __name__ == "__main__":
    main()
